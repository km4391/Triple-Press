#use python  to make a MS word document an excel sspradsheet and a powerpoint presentiation slide from this script's total output.




from ibapi.client import EClient
from ibapi.wrapper import EWrapper
from ibapi.contract import Contract
from threading import Timer 
import pandas as pd
import xml.etree.ElementTree as ET

class TestApp(EWrapper, EClient):
    def __init__(self):
        EClient.__init__(self, self)
        #stores data
        self.data = []
 
    def error(self, reqId, errorCode, errorString,caster):
        print("Error: ", reqId, " ", errorCode, " ", errorString," ",caster)
 
    def nextValidId(self, orderId):
        self.start()
 
    def updatePortfolio(self, contract: Contract, position: float, marketPrice: float, marketValue: float,
                        averageCost: float, unrealizedPNL: float, realizedPNL: float, accountName: str):
        row = {
            'symbol': contract.symbol, 
            'secType': contract.secType, 
            'exchange': contract.exchange, 
            'position': position, 
            'marketPrice': marketPrice, 
            'marketValue': marketValue, 
            'averageCost': averageCost, 
            'unrealizedPNL': unrealizedPNL, 
            'realizedPNL': realizedPNL, 
            'accountName': accountName
        }
        
        self.data.append(row)
 
    def updateAccountValue(self, key: str, val: str, currency: str, accountName: str):
        print("UpdateAccountValue. Key:", key, "Value:", val, "Currency:", currency, "AccountName:", accountName)
 
    def updateAccountTime(self, timeStamp: str):
        print("UpdateAccountTime. Time:", timeStamp)
 
    def accountDownloadEnd(self, accountName: str):
        print("AccountDownloadEnd. Account:", accountName)
        #store for display
        self.data_df = pd.DataFrame(self.data)
        self.display()
 
    def start(self):
        # Account number can be omitted when using reqAccountUpdates with single account structure
        self.reqAccountUpdates(True, "")
 
    def stop(self):
        self.reqAccountUpdates(False, "")
        self.done = True
        self.disconnect()
        #display table with updated values
        self.display()
        
    def display(self):
        #call and print out the  dataframe
        print(self.data_df)
 
        #Create word document using xml
        root = ET.Element('word')
        doc = ET.SubElement(root, 'document')
        body = ET.SubElement(doc, 'body')
        table = ET.SubElement(body, 'table')
        row = ET.SubElement(table, 'row')
        for col in self.data_df.columns:
            cell = ET.SubElement(row, 'cell')
            text = ET.SubElement(cell, 'text')
            text.text = col
        for index, row in self.data_df.iterrows():
            new_row = ET.SubElement(table, 'row')
            for col in self.data_df.columns:
                cell = ET.SubElement(new_row, 'cell')
                text = ET.SubElement(cell, 'text')
                text.text = str(row[col])
        tree = ET.ElementTree(root)
        tree.write('word_doc.xml')
        root = ET.Element('word')
        doc = ET.SubElement(root, 'document')
        body = ET.SubElement(doc, 'body')
        table = ET.SubElement(body, 'table')
        row = ET.SubElement(table, 'row')
        for col in self.data_df.columns:
            cell = ET.SubElement(row, 'cell')
            text = ET.SubElement(cell, 'text')
            text.text = col
        for index, row in self.data_df.iterrows():
            new_row = ET.SubElement(table, 'row')
            for col in self.data_df.columns:
                cell = ET.SubElement(new_row, 'cell')
                text = ET.SubElement(cell, 'text')
                text.text = str(row[col])
        tree = ET.ElementTree(root)
        tree.write('word_doc.xml')

        # Create excel spreadsheet 
        self.data_df.to_excel('spreadsheet.xlsx')

        #Create powerpoint presentation slide
        import pptx
        from pptx.util import Inches

        # Create a new presentation
        prs = pptx.Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = 'Data Output'

        # Add table
        rows = cols = len(self.data_df.columns)
        top = Inches(1.5)
        left = Inches(3.5)
        width = Inches(7.0)
        height = Inches(4.5)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column labels
        for col in self.data_df.columns:
            table.cell(0, self.data_df.columns.get_loc(col)).text = col

        # Add rows
        for index, row in self.data_df.iterrows():
            for col in self.data_df.columns:
                table.cell(index+1, self.data_df.columns.get_loc(col)).text = str(row[col])

        # Save
        prs.save('presentation.pptx')
 
def main():
    app = TestApp()
    app.connect("127.0.0.1", 7495, 0)
 
    Timer(5, app.stop).start() 
    
    app.run()
 
if __name__ == "__main__":
    main()


#Here rewrite please:

# Create word document using xml
